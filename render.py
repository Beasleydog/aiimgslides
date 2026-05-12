import math

from PIL import Image, ImageDraw
from pptx.oxml.ns import qn
from pptx.oxml.xmlchemy import OxmlElement
from pptx import Presentation
from pptx.util import Inches

import config
from chart_element import add_chart_to_png, add_chart_to_pptx
from connector_element import add_connector_to_png, add_connector_to_pptx
from freeform_element import add_freeform_to_png, add_freeform_to_pptx
from image_element import add_image_to_png, add_image_to_pptx
from shape_element import add_shape_to_png, add_shape_to_pptx
from table_element import add_table_to_png, add_table_to_pptx
from text_element import add_text_to_png, add_text_to_pptx
from text_fit import fit_text_elements
from svg_element import add_svg_image_to_png, add_svg_to_png
from utils import ppt_color, px_box


PATH_GRADIENT_FOCUS_RECTS = {
    "center": {"l": "50000", "t": "50000", "r": "50000", "b": "50000"},
    "top_left": {"l": "100000", "t": "100000", "r": "0", "b": "0"},
    "top_right": {"l": "0", "t": "100000", "r": "100000", "b": "0"},
    "bottom_left": {"l": "100000", "t": "0", "r": "0", "b": "100000"},
    "bottom_right": {"l": "0", "t": "0", "r": "100000", "b": "100000"},
}


def set_path_gradient(shape, gradient_type, focus):
    grad_fill = shape._element.spPr.find(qn("a:gradFill"))
    if grad_fill is None:
        return
    for child_name in ("a:lin", "a:path"):
        child = grad_fill.find(qn(child_name))
        if child is not None:
            grad_fill.remove(child)

    path = OxmlElement("a:path")
    path.set("path", "rect" if gradient_type == "rectangular" else "shape" if gradient_type == "shape" else "circle")
    fill_to_rect = OxmlElement("a:fillToRect")
    for attr, value in PATH_GRADIENT_FOCUS_RECTS.get(focus, PATH_GRADIENT_FOCUS_RECTS["center"]).items():
        fill_to_rect.set(attr, value)
    path.append(fill_to_rect)
    grad_fill.append(path)


def add_background_to_pptx(slide, el, image_path):
    style = el.data.get("style", "solid")
    shape = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(config.SLIDE_W), Inches(config.SLIDE_H))
    if style == "gradient":
        shape.fill.gradient()
        gradient_type = el.data.get("gradient_type", "linear")
        if gradient_type == "linear":
            shape.fill.gradient_angle = el.data.get("angle", 0)
        else:
            set_path_gradient(shape, gradient_type, el.data.get("focus", "center"))
        colors = el.data.get("colors", [el.data["fill"], el.data["fill"]])
        for idx, stop in enumerate(shape.fill.gradient_stops):
            stop.color.rgb = ppt_color(colors[min(idx, len(colors) - 1)])
    elif style == "pattern":
        shape.fill.patterned()
        shape.fill.pattern = el.data["pattern"]
        shape.fill.fore_color.rgb = ppt_color(el.data["fore_color"])
        shape.fill.back_color.rgb = ppt_color(el.data["back_color"])
    else:
        shape.fill.solid()
        shape.fill.fore_color.rgb = ppt_color(el.data["fill"])

    shape.line.color.rgb = ppt_color(el.data.get("fill", (255, 255, 255)))


def gradient_image(size, colors, angle=0, gradient_type="linear", focus="center"):
    w, h = size
    c1, c2 = colors[0], colors[-1]
    img = Image.new("RGB", size, c1)
    pixels = img.load()
    focus_points = {
        "center": (w / 2, h / 2),
        "top_left": (0, 0),
        "top_right": (w, 0),
        "bottom_left": (0, h),
        "bottom_right": (w, h),
    }
    fx, fy = focus_points.get(focus, focus_points["center"])
    if gradient_type == "linear":
        radians = angle / 180 * math.pi
        dx = math.cos(radians)
        dy = math.sin(radians)
        values = [x * dx + y * dy for x, y in [(0, 0), (w, 0), (0, h), (w, h)]]
        lo, hi = min(values), max(values)
        span = max(1, hi - lo)
    else:
        max_dist = max(math.hypot(x - fx, y - fy) for x, y in [(0, 0), (w, 0), (0, h), (w, h)])
    for y in range(h):
        for x in range(w):
            if gradient_type == "linear":
                t = ((x * dx + y * dy) - lo) / span
            elif gradient_type == "rectangular":
                t = max(abs(x - fx) / max(1, w), abs(y - fy) / max(1, h)) * 2
            else:
                t = math.hypot(x - fx, y - fy) / max(1, max_dist)
            t = max(0, min(1, t))
            pixels[x, y] = tuple(int(a * (1 - t) + b * t) for a, b in zip(c1, c2))
    return img


def pattern_image(size, fore, back, pattern_name):
    img = Image.new("RGB", size, back)
    draw = ImageDraw.Draw(img)
    w, h = size
    name = str(pattern_name).lower()
    step = 18
    if "grid" in name or "cross" in name or "trellis" in name:
        for x in range(0, w, step):
            draw.line((x, 0, x, h), fill=fore, width=1)
        for y in range(0, h, step):
            draw.line((0, y, w, y), fill=fore, width=1)
    elif "diagonal" in name or "zig_zag" in name:
        for x in range(-h, w, step):
            draw.line((x, h, x + h, 0), fill=fore, width=2)
    elif "horizontal" in name or "wave" in name:
        for y in range(0, h, step):
            draw.line((0, y, w, y), fill=fore, width=2)
    elif "vertical" in name:
        for x in range(0, w, step):
            draw.line((x, 0, x, h), fill=fore, width=2)
    else:
        for y in range(0, h, step):
            for x in range(0, w, step):
                draw.ellipse((x, y, x + 3, y + 3), fill=fore)
    return img


def add_background_to_png(canvas, image, el):
    style = el.data.get("style", "solid")
    if style == "gradient":
        canvas.paste(
            gradient_image(
                canvas.size,
                el.data.get("colors", [el.data["fill"], el.data["fill"]]),
                el.data.get("angle", 0),
                el.data.get("gradient_type", "linear"),
                el.data.get("focus", "center"),
            )
        )
    elif style == "pattern":
        canvas.paste(pattern_image(canvas.size, el.data["fore_color"], el.data["back_color"], el.data["pattern"]))
    else:
        canvas.paste(el.data["fill"], [0, 0, config.PREVIEW_W, config.PREVIEW_H])


def add_elements_to_slide(slide, elements, image_path):
    for el in elements:
        if el.kind == "background":
            add_background_to_pptx(slide, el, image_path)
        elif el.kind == "text":
            add_text_to_pptx(slide, el)
        elif el.kind == "shape":
            add_shape_to_pptx(slide, el)
        elif el.kind == "table":
            add_table_to_pptx(slide, el)
        elif el.kind == "image":
            add_image_to_pptx(slide, el, image_path)
        elif el.kind == "connector":
            add_connector_to_pptx(slide, el)
        elif el.kind == "chart":
            add_chart_to_pptx(slide, el)
        elif el.kind == "freeform":
            add_freeform_to_pptx(slide, el)
        elif el.kind == "svg":
            pass
        elif el.kind == "svg_image":
            pass


def add_to_pptx(elements, image_path):
    elements = fit_text_elements(elements)
    prs = Presentation()
    prs.slide_width = Inches(config.SLIDE_W)
    prs.slide_height = Inches(config.SLIDE_H)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_elements_to_slide(slide, elements, image_path)
    prs.save(config.OUTPUT_PPTX)


def add_single_slide_to_pptx(elements, image_path, output_pptx):
    elements = fit_text_elements(elements)
    prs = Presentation()
    prs.slide_width = Inches(config.SLIDE_W)
    prs.slide_height = Inches(config.SLIDE_H)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_elements_to_slide(slide, elements, image_path)
    prs.save(output_pptx)


def add_deck_to_pptx(slide_elements, image_path, output_pptx):
    slide_elements = [fit_text_elements(elements) for elements in slide_elements]
    prs = Presentation()
    prs.slide_width = Inches(config.SLIDE_W)
    prs.slide_height = Inches(config.SLIDE_H)
    for elements in slide_elements:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        add_elements_to_slide(slide, elements, image_path)
    prs.save(output_pptx)


def add_to_png(elements, image):
    elements = fit_text_elements(elements)
    canvas = Image.new("RGB", (config.PREVIEW_W, config.PREVIEW_H), config.PREVIEW_BACKGROUND)
    draw = ImageDraw.Draw(canvas, "RGBA")

    for el in elements:
        box = px_box(el)
        if el.kind == "background":
            add_background_to_png(canvas, image, el)
        elif el.kind == "text":
            add_text_to_png(draw, el, box)
        elif el.kind == "shape":
            add_shape_to_png(draw, el, box)
        elif el.kind == "table":
            add_table_to_png(draw, el, box)
        elif el.kind == "image":
            add_image_to_png(canvas, image, box)
        elif el.kind == "connector":
            add_connector_to_png(draw, el, box)
        elif el.kind == "chart":
            add_chart_to_png(draw, el, box)
        elif el.kind == "freeform":
            add_freeform_to_png(draw, el, box)
        elif el.kind == "svg":
            add_svg_to_png(draw, el, box)
        elif el.kind == "svg_image":
            add_svg_image_to_png(draw, el, box)

    canvas.save(config.OUTPUT_PNG)
