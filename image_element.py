from io import BytesIO
import random

import requests
from PIL import Image, ImageDraw
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.util import Inches, Pt

import config
from element_model import Element
from utils import ppt_color, rand_color


def make_image(x, y, w, h):
    crop = random_crop() if random.random() < config.IMAGE_CROP_PROBABILITY else (0, 0, 0, 0)
    return Element(
        "image",
        x,
        y,
        w,
        h,
        {
            "url": config.EXAMPLE_IMAGE_URL,
            "crop": crop,
            "mask_overlay": random.random() < config.IMAGE_MASK_OVERLAY_PROBABILITY,
            "mask_line": rand_color(),
        },
    )


def random_crop():
    max_crop = config.IMAGE_MAX_CROP
    left = random.uniform(0, max_crop)
    right = random.uniform(0, max_crop)
    top = random.uniform(0, max_crop)
    bottom = random.uniform(0, max_crop)
    return left, top, right, bottom


def download_image():
    try:
        res = requests.get(config.EXAMPLE_IMAGE_URL, timeout=15)
        res.raise_for_status()
        return Image.open(BytesIO(res.content)).convert("RGB")
    except Exception:
        img = Image.new("RGB", (900, 600), (230, 232, 236))
        draw = ImageDraw.Draw(img)
        draw.rectangle((60, 60, 840, 540), outline=(90, 100, 115), width=8)
        draw.text((330, 275), "example image", fill=(60, 70, 85))
        return img


def add_image_to_pptx(slide, el, image_path):
    pic = slide.shapes.add_picture(str(image_path), Inches(el.x), Inches(el.y), Inches(el.w), Inches(el.h))
    left, top, right, bottom = el.data.get("crop", (0, 0, 0, 0))
    pic.crop_left = left
    pic.crop_top = top
    pic.crop_right = right
    pic.crop_bottom = bottom
    if el.data["mask_overlay"]:
        shape = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
            Inches(el.x),
            Inches(el.y),
            Inches(el.w),
            Inches(el.h),
        )
        shape.fill.background()
        shape.line.color.rgb = ppt_color(el.data["mask_line"])
        shape.line.width = Pt(1.2)


def add_image_to_png(canvas, image, box):
    crop = image.copy()
    crop.thumbnail((box[2] - box[0], box[3] - box[1]))
    fitted = Image.new("RGB", (box[2] - box[0], box[3] - box[1]), (230, 230, 230))
    fitted.paste(crop, ((fitted.width - crop.width) // 2, (fitted.height - crop.height) // 2))
    canvas.paste(fitted, box[:2])
