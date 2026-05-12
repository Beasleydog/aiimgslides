import argparse
import json
import re
from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.util import Inches, Pt

from compact_schema import compact_to_full
from train import MODEL_NAME, OUTPUT_DIR, USER_PROMPT, balanced_json_slice
from train import IMAGE_MAX_PIXELS, IMAGE_MIN_PIXELS


DEFAULT_IMAGE = Path("realslide.png")
DEFAULT_OUT_DIR = Path("val_output")
SLIDE_W = 13.333
SLIDE_H = 7.5


SHAPE_TYPES = {
    "rect": MSO_AUTO_SHAPE_TYPE.RECTANGLE,
    "rectangle": MSO_AUTO_SHAPE_TYPE.RECTANGLE,
    "round_rect": MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
    "rounded_rectangle": MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
    "ellipse": MSO_AUTO_SHAPE_TYPE.OVAL,
    "oval": MSO_AUTO_SHAPE_TYPE.OVAL,
    "circle": MSO_AUTO_SHAPE_TYPE.OVAL,
    "triangle": MSO_AUTO_SHAPE_TYPE.ISOSCELES_TRIANGLE,
    "right_triangle": MSO_AUTO_SHAPE_TYPE.RIGHT_TRIANGLE,
    "diamond": MSO_AUTO_SHAPE_TYPE.DIAMOND,
    "right_arrow": MSO_AUTO_SHAPE_TYPE.RIGHT_ARROW,
    "left_arrow": MSO_AUTO_SHAPE_TYPE.LEFT_ARROW,
    "up_arrow": MSO_AUTO_SHAPE_TYPE.UP_ARROW,
    "down_arrow": MSO_AUTO_SHAPE_TYPE.DOWN_ARROW,
    "line": MSO_AUTO_SHAPE_TYPE.RECTANGLE,
}

CHART_TYPES = {
    "bar": XL_CHART_TYPE.BAR_CLUSTERED,
    "column": XL_CHART_TYPE.COLUMN_CLUSTERED,
    "line": XL_CHART_TYPE.LINE_MARKERS,
    "area": XL_CHART_TYPE.AREA,
    "pie": XL_CHART_TYPE.PIE,
    "doughnut": XL_CHART_TYPE.DOUGHNUT,
}


def color(value, fallback=(30, 30, 30)):
    if isinstance(value, str):
        text = value.strip().lstrip("#")
        if re.fullmatch(r"[0-9a-fA-F]{6}", text):
            return RGBColor(int(text[0:2], 16), int(text[2:4], 16), int(text[4:6], 16))
    if isinstance(value, dict):
        value = value.get("rgb") or value.get("value") or value.get("color")
    if isinstance(value, (list, tuple)) and len(value) >= 3:
        try:
            return RGBColor(*(max(0, min(255, int(v))) for v in value[:3]))
        except Exception:
            pass
    return RGBColor(*fallback)


def num(value, fallback):
    try:
        return float(value)
    except Exception:
        return fallback


def bbox(obj):
    box = obj.get("bbox") if isinstance(obj, dict) else {}
    x = max(0, min(SLIDE_W, num(box.get("x"), 0.5)))
    y = max(0, min(SLIDE_H, num(box.get("y"), 0.5)))
    w = max(0.1, min(SLIDE_W - x, num(box.get("w"), 2.0)))
    h = max(0.1, min(SLIDE_H - y, num(box.get("h"), 1.0)))
    return Inches(x), Inches(y), Inches(w), Inches(h)


def props(obj):
    data = obj.get("properties", {}) if isinstance(obj, dict) else {}
    return data if isinstance(data, dict) else {}


def extract_scene(text):
    match = re.search(r"<json>\s*(.*?)\s*</json>", text, flags=re.DOTALL | re.IGNORECASE)
    raw = match.group(1) if match else balanced_json_slice(text)
    return json.loads(raw)


def load_model(model_dir, base_model):
    import torch
    from peft import PeftModel
    from transformers import AutoProcessor, Qwen2_5_VLForConditionalGeneration

    dtype = torch.bfloat16 if torch.cuda.is_available() and torch.cuda.is_bf16_supported() else torch.float16
    model = Qwen2_5_VLForConditionalGeneration.from_pretrained(
        base_model,
        device_map="auto",
        dtype=dtype,
        attn_implementation="eager",
    )
    model.config.use_cache = True
    if (model_dir / "adapter_config.json").exists():
        model = PeftModel.from_pretrained(model, model_dir)
    processor_source = model_dir if (model_dir / "preprocessor_config.json").exists() else base_model
    processor = AutoProcessor.from_pretrained(
        processor_source,
        use_fast=True,
        min_pixels=IMAGE_MIN_PIXELS,
        max_pixels=IMAGE_MAX_PIXELS,
    )
    if hasattr(processor, "tokenizer"):
        processor.tokenizer.model_max_length = 131072
    model.eval()
    return model, processor


def generate_scene(model, processor, image_path, max_new_tokens):
    import torch

    image = Image.open(image_path).convert("RGB")
    messages = [
        {
            "role": "user",
            "content": [
                {"type": "image", "image": image},
                {"type": "text", "text": USER_PROMPT},
            ],
        }
    ]
    prompt = processor.apply_chat_template(messages, tokenize=False, add_generation_prompt=True)
    inputs = processor(text=[prompt], images=[image], return_tensors="pt", padding=True, truncation=False)
    device = next(model.parameters()).device
    inputs = {key: value.to(device) if hasattr(value, "to") else value for key, value in inputs.items()}
    with torch.inference_mode():
        output_ids = model.generate(
            **inputs,
            max_new_tokens=max_new_tokens,
            do_sample=False,
            remove_invalid_values=True,
            renormalize_logits=True,
        )
    generated = output_ids[:, inputs["input_ids"].shape[1] :]
    text = processor.batch_decode(generated, skip_special_tokens=True)[0]
    return text, compact_to_full(extract_scene(text), image_file=str(image_path))


def add_text(slide, obj):
    data = props(obj)
    shape = slide.shapes.add_textbox(*bbox(obj))
    frame = shape.text_frame
    frame.word_wrap = True
    frame.clear()
    p = frame.paragraphs[0]
    run = p.add_run()
    run.text = str(data.get("text") or data.get("content") or data.get("label") or "")
    run.font.size = Pt(max(6, min(72, int(num(data.get("font_size"), 18)))))
    run.font.name = str(data.get("font_face") or data.get("font") or "Aptos")
    run.font.bold = bool(data.get("bold", False))
    run.font.italic = bool(data.get("italic", False))
    run.font.underline = bool(data.get("underline", False))
    run.font.color.rgb = color(data.get("color") or data.get("text_color"), (20, 20, 20))
    if data.get("bg_color") or data.get("background"):
        shape.fill.solid()
        shape.fill.fore_color.rgb = color(data.get("bg_color") or data.get("background"), (255, 255, 255))


def add_shape(slide, obj):
    data = props(obj)
    name = str(data.get("shape") or data.get("kind") or "rect").lower()
    shape_type = SHAPE_TYPES.get(name, MSO_AUTO_SHAPE_TYPE.RECTANGLE)
    shape = slide.shapes.add_shape(shape_type, *bbox(obj))
    if bool(data.get("outline_only", False)):
        shape.fill.background()
    else:
        shape.fill.solid()
        shape.fill.fore_color.rgb = color(data.get("fill") or data.get("fill_color"), (230, 230, 230))
    shape.line.color.rgb = color(data.get("line") or data.get("line_color") or data.get("border"), (40, 40, 40))
    shape.line.width = Pt(max(0, min(12, num(data.get("line_width") or data.get("border_width"), 1.0))))


def add_table(slide, obj):
    data = props(obj)
    rows = max(1, min(12, int(num(data.get("rows"), 3))))
    cols = max(1, min(8, int(num(data.get("cols"), 3))))
    table = slide.shapes.add_table(rows, cols, *bbox(obj)).table
    cells = data.get("cells") if isinstance(data.get("cells"), list) else []
    for r in range(rows):
        for c in range(cols):
            cell = table.cell(r, c)
            value = cells[r][c] if r < len(cells) and isinstance(cells[r], list) and c < len(cells[r]) else ""
            cell.text = str(value)
            cell.fill.solid()
            cell.fill.fore_color.rgb = color(data.get("header") if r == 0 else data.get("body"), (245, 245, 245))


def add_chart(slide, obj):
    data = props(obj)
    chart_type = CHART_TYPES.get(str(data.get("chart") or "column").lower(), XL_CHART_TYPE.COLUMN_CLUSTERED)
    categories = data.get("categories") if isinstance(data.get("categories"), list) else ["A", "B", "C"]
    series = data.get("series") if isinstance(data.get("series"), list) else []
    chart_data = CategoryChartData()
    chart_data.categories = [str(item) for item in categories[:8]]
    if not series:
        series = [{"name": "Series", "values": [1 for _ in chart_data.categories]}]
    for item in series[:4]:
        values = item.get("values") if isinstance(item, dict) and isinstance(item.get("values"), list) else []
        parsed_values = [num(v, 1) for v in values[: len(categories)]]
        while len(parsed_values) < len(categories):
            parsed_values.append(1)
        chart_data.add_series(str(item.get("name", "Series")), parsed_values)
    slide.shapes.add_chart(chart_type, *bbox(obj), chart_data)


def add_background(slide, scene):
    bg = scene.get("background") if isinstance(scene, dict) else None
    data = props(bg) if bg else {}
    fill = data.get("fill") or data.get("color") or data.get("back_color") or (255, 255, 255)
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0), Inches(0), Inches(SLIDE_W), Inches(SLIDE_H))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color(fill, (255, 255, 255))
    shape.line.fill.background()


def scene_to_pptx(scene, image_path, output_pptx):
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, scene)

    objects = scene.get("objects") if isinstance(scene, dict) else []
    if not isinstance(objects, list):
        objects = []
    objects = sorted(objects, key=lambda item: num(item.get("z_order") if isinstance(item, dict) else 0, 0))
    for obj in objects:
        kind = str(obj.get("type", "")).lower() if isinstance(obj, dict) else ""
        try:
            if kind == "text":
                add_text(slide, obj)
            elif kind in {"shape", "freeform", "svg"}:
                add_shape(slide, obj)
            elif kind == "table":
                add_table(slide, obj)
            elif kind == "chart":
                add_chart(slide, obj)
            elif kind in {"image", "svg_image"}:
                slide.shapes.add_picture(str(image_path), *bbox(obj))
        except Exception as exc:
            print(f"Skipped object {obj.get('id', '?')}: {exc}")

    prs.save(output_pptx)


def main():
    parser = argparse.ArgumentParser()
    parser.add_argument("--image", type=Path, default=DEFAULT_IMAGE)
    parser.add_argument("--model-dir", type=Path, default=OUTPUT_DIR)
    parser.add_argument("--base-model", default=MODEL_NAME)
    parser.add_argument("--out-dir", type=Path, default=DEFAULT_OUT_DIR)
    parser.add_argument("--max-new-tokens", type=int, default=2048)
    args = parser.parse_args()

    args.out_dir.mkdir(parents=True, exist_ok=True)
    raw_path = args.out_dir / "val_raw.txt"
    json_path = args.out_dir / "val.json"
    pptx_path = args.out_dir / "val.pptx"

    model, processor = load_model(args.model_dir, args.base_model)
    raw_text, scene = generate_scene(model, processor, args.image, args.max_new_tokens)
    raw_path.write_text(raw_text, encoding="utf-8")
    json_path.write_text(json.dumps(scene, indent=2, sort_keys=True), encoding="utf-8")
    scene_to_pptx(scene, args.image, pptx_path)
    print(f"Wrote {pptx_path}")
    print(f"Wrote {json_path}")
    print(f"Wrote {raw_path}")


if __name__ == "__main__":
    main()
