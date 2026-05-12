import random

from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.util import Inches, Pt

from element_model import Element
from utils import add_wrapped_text, rand_color, rand_text


CHART_TYPES = [
    ("column", XL_CHART_TYPE.COLUMN_CLUSTERED),
    ("bar", XL_CHART_TYPE.BAR_CLUSTERED),
    ("line", XL_CHART_TYPE.LINE_MARKERS),
    ("area", XL_CHART_TYPE.AREA),
    ("pie", XL_CHART_TYPE.PIE),
    ("doughnut", XL_CHART_TYPE.DOUGHNUT),
]


def make_chart(x, y, w, h):
    name, chart_type = random.choice(CHART_TYPES)
    categories = [rand_text(1, 1) for _ in range(random.randint(3, 6))]
    series_count = 1 if name in {"pie", "doughnut"} else random.randint(1, 3)
    return Element(
        "chart",
        x,
        y,
        w,
        h,
        {
            "chart": name,
            "ppt_type": chart_type,
            "categories": categories,
            "series": [
                {
                    "name": rand_text(1, 2),
                    "values": [random.randint(8, 95) for _ in categories],
                    "color": rand_color(),
                }
                for _ in range(series_count)
            ],
            "has_legend": random.choice([True, False]),
            "has_title": random.choice([True, False]),
            "title": rand_text(2, 5),
        },
    )


def add_chart_to_pptx(slide, el):
    data = CategoryChartData()
    data.categories = el.data["categories"]
    for series in el.data["series"]:
        data.add_series(series["name"], series["values"])

    frame = slide.shapes.add_chart(
        el.data["ppt_type"],
        Inches(el.x),
        Inches(el.y),
        Inches(el.w),
        Inches(el.h),
        data,
    )
    chart = frame.chart
    chart.has_legend = el.data["has_legend"]
    chart.has_title = el.data["has_title"]
    if chart.has_title:
        chart.chart_title.text_frame.text = el.data["title"]
        chart.chart_title.text_frame.paragraphs[0].runs[0].font.size = Pt(10)


def add_chart_to_png(draw, el, box):
    x1, y1, x2, y2 = box
    draw.rectangle(box, fill=(255, 255, 255, 225), outline=(50, 50, 50), width=2)
    if el.data["has_title"]:
        add_wrapped_text(draw, el.data["title"], (x1 + 6, y1 + 5, x2 - 6, y1 + 30), (35, 35, 35), 14, True)
    values = el.data["series"][0]["values"]
    max_value = max(values)
    plot_top = y1 + 34
    plot_bottom = y2 - 18
    slot = max(1, (x2 - x1 - 20) / len(values))
    for i, value in enumerate(values):
        bx1 = int(x1 + 10 + i * slot + slot * 0.18)
        bx2 = int(x1 + 10 + (i + 1) * slot - slot * 0.18)
        bh = int((plot_bottom - plot_top) * value / max_value)
        draw.rectangle((bx1, plot_bottom - bh, bx2, plot_bottom), fill=(*el.data["series"][0]["color"], 180))
