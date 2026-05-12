import random

from pptx.enum.text import MSO_AUTO_SIZE
from pptx.util import Inches, Pt

import config
from element_model import Element
from utils import add_wrapped_text, ppt_color, rand_color, rand_text, weighted_choice


FONT_FACES = ["Aptos", "Calibri", "Arial", "Georgia", "Trebuchet MS", "Verdana"]
ALIGNMENTS = [None, 1, 2, 3]


def text_word_range(kind):
    ranges = {
        "kicker": (1, 3),
        "header": (2, 7),
        "subhead": (5, 14),
        "body": (18, 45),
        "caption": (4, 12),
    }
    return ranges[kind]


def make_text(x, y, w, h, forced_kind=None):
    kind = forced_kind or weighted_choice(config.TEXT_KIND_WEIGHTS)
    min_words, max_words = text_word_range(kind)
    min_size, max_size = config.TEXT_FONT_SIZES[kind]
    bg_color = rand_color() if random.random() < 0.18 else None
    return Element(
        "text",
        x,
        y,
        w,
        h,
        {
            "kind": kind,
            "text": rand_text(min_words, max_words),
            "font_size": random.randint(min_size, max_size),
            "font_face": random.choice(FONT_FACES),
            "color": rand_color(),
            "bold": kind == "header" or random.random() < 0.3,
            "italic": random.random() < 0.16,
            "underline": random.random() < 0.08,
            "bg_color": bg_color,
            "margin": random.uniform(0.02, 0.10),
        },
    )


def add_text_to_pptx(slide, el):
    box = slide.shapes.add_textbox(Inches(el.x), Inches(el.y), Inches(el.w), Inches(el.h))
    if el.data["bg_color"]:
        box.fill.solid()
        box.fill.fore_color.rgb = ppt_color(el.data["bg_color"])
    box.text_frame.word_wrap = True
    box.text_frame.auto_size = MSO_AUTO_SIZE.NONE
    p = box.text_frame.paragraphs[0]
    run = p.add_run()
    run.text = el.data["text"]
    run.font.size = Pt(el.data["font_size"])
    run.font.name = el.data["font_face"]
    run.font.bold = el.data["bold"]
    run.font.italic = el.data["italic"]
    run.font.underline = el.data["underline"]
    run.font.color.rgb = ppt_color(el.data["color"])
    margin = Inches(el.data["margin"])
    box.text_frame.margin_left = margin
    box.text_frame.margin_right = margin
    box.text_frame.margin_top = margin
    box.text_frame.margin_bottom = margin


def add_text_to_png(draw, el, box):
    if el.data["bg_color"]:
        draw.rounded_rectangle(box, radius=4, fill=(*el.data["bg_color"], 115))
    size = max(10, int(el.data["font_size"] / 72 * (config.PREVIEW_W / config.SLIDE_W)))
    x1, y1, x2, y2 = box
    pad = int(el.data["margin"] / config.SLIDE_W * config.PREVIEW_W)
    add_wrapped_text(draw, el.data["text"], (x1 + pad, y1 + pad, x2 - pad, y2 - pad), el.data["color"], size, el.data["bold"])
